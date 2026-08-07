"""Builds the presentation deck for the Пробок.Нет pack.

    python tools/docs-media/build-deck.py <pack-dir>

Reads the screenshots the pack already contains and writes the .pptx beside
them, so the deck can be rebuilt from a fresh capture without touching slides
by hand. The look follows the product's dark theme rather than a stock
template: the same yellow, the same paper, the same three flow colours.
"""
from __future__ import annotations

import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

PACK = Path(sys.argv[1] if len(sys.argv) > 1 else "tmp/FOR_INVESTORS")
SHOTS = PACK / "screenshots"
BRAND = Path("apps/web/public/brand/logo-dark-512.png")

INK = RGBColor(0xF1, 0xF0, 0xEE)
MUTED = RGBColor(0x9A, 0x98, 0x94)
PAPER = RGBColor(0x1B, 0x1B, 0x1E)
PAPER_2 = RGBColor(0x24, 0x24, 0x28)
ACCENT = RGBColor(0xFF, 0xD1, 0x1A)
GREEN = RGBColor(0x3F, 0xC4, 0x7B)
RED = RGBColor(0xFF, 0x5D, 0x5D)
FONT = "Segoe UI"

W, H = Inches(13.333), Inches(7.5)

deck = Presentation()
deck.slide_width, deck.slide_height = W, H
BLANK = deck.slide_layouts[6]


def slide(dark: bool = True):
    page = deck.slides.add_slide(BLANK)
    fill = page.background.fill
    fill.solid()
    fill.fore_color.rgb = PAPER if dark else RGBColor(0xFF, 0xFF, 0xFF)
    return page


def text(page, body, left, top, width, height, *, size=18, colour=INK, bold=False,
         align=PP_ALIGN.LEFT, spacing=1.25, italic=False):
    box = page.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    lines = body if isinstance(body, list) else [body]
    for index, line in enumerate(lines):
        para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        para.alignment = align
        para.line_spacing = spacing
        para.space_after = Pt(6)
        run = para.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = colour
        run.font.name = FONT
    return box


def rule(page, top, *, left=Inches(0.9), width=Inches(1.6), colour=ACCENT, thickness=Pt(3)):
    line = page.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = colour
    line.line.width = thickness


def heading(page, title, *, top=Inches(0.55), colour=INK, size=34, underline=True):
    """Title plus its accent rule, placed below however many lines the title took.

    A fixed rule position struck a two-line heading through the middle, which is
    exactly what it looked like: a deletion.
    """
    # Roughly how much of the 11.5in line a bold Segoe UI character claims.
    per_line = int(11.5 * 72 / (size * 0.52))
    lines = 1 + (len(title) - 1) // per_line
    text(page, title, Inches(0.9), top, W - Inches(1.8), Inches(1.0) * lines,
         size=size, colour=colour, bold=True, spacing=1.05)
    if underline:
        rule(page, top + Inches(0.62) * lines + Inches(0.38), colour=colour if colour is RED else ACCENT)
    return top + Inches(0.62) * lines + Inches(0.75)


def picture(page, name, box_left, box_top, box_width, box_height, *, shadow=True):
    """Fits an image inside a box, centred, without distorting it."""
    path = Path(name)
    if not path.is_absolute():
        path = SHOTS / name if (SHOTS / name).exists() else PACK / name
    with Image.open(path) as image:
        ratio = image.width / image.height
    if box_width / box_height > ratio:
        height, width = box_height, Emu(int(box_height * ratio))
    else:
        width, height = box_width, Emu(int(box_width / ratio))
    left = box_left + Emu(int((box_width - width) / 2))
    top = box_top + Emu(int((box_height - height) / 2))
    shape = page.shapes.add_picture(str(path), left, top, width, height)
    if not shadow:
        shape.shadow.inherit = False
    return shape


def caption(page, body, top):
    text(page, body, Inches(0.9), top, W - Inches(1.8), Inches(0.5),
         size=12, colour=MUTED, italic=True)


def bullets(page, items, left, top, width, height, *, size=17):
    return text(page, [f"•  {item}" for item in items], left, top, width, height,
                size=size, colour=INK, spacing=1.35)


def stat(page, value, label, left, top, width):
    text(page, value, left, top, width, Inches(0.9), size=40, colour=ACCENT, bold=True, spacing=1.0)
    # Split rather than embed a newline in a run: a run is one line, and the
    # break would be swallowed.
    text(page, label.splitlines(), left, top + Inches(0.85), width, Inches(1.1),
         size=14, colour=MUTED, spacing=1.2)


# ─────────────────────────────────────────────────────────── 1. Title
page = slide()
if BRAND.exists():
    page.shapes.add_picture(str(BRAND), Inches(0.9), Inches(1.5), Inches(1.1), Inches(1.1))
text(page, "Пробок.Нет", Inches(0.9), Inches(2.8), Inches(10), Inches(1.3),
     size=60, colour=INK, bold=True, spacing=1.0)
text(page, "Маршрут не самый быстрый, а самый зелёный", Inches(0.9), Inches(4.0),
     Inches(11), Inches(0.8), size=26, colour=ACCENT, spacing=1.1)
text(page, ["Навигаторы отвечают на вопрос «как доехать быстрее».",
            "Пробок.Нет отвечает на другой: «как доехать, почти не стоя в пробке»."],
     Inches(0.9), Inches(4.9), Inches(11), Inches(1.4), size=17, colour=MUTED)

# ─────────────────────────────────────────────────────────── 2. Problem
page = slide()
heading(page, "Навигатор считает минуты. Человек — нет")
text(page,
     ["Любой навигатор оптимизирует ETA. Ему всё равно, из чего сложилось время: "
      "50 минут свободного хода и 50 минут в стоящем потоке для него одинаковы.",
      "Поэтому он отправит вас в плотный поток на магистрали ради выигрыша в две минуты — "
      "и вы будете стоять 40 минут из 50."],
     Inches(0.9), Inches(2.0), Inches(11.5), Inches(2.0), size=20)
text(page,
     "Стоять в пробке — отдельная цена: усталость, расход, сорванные звонки, "
     "невозможность спланировать приезд. Многие готовы ехать дольше, лишь бы ехать.",
     Inches(0.9), Inches(4.2), Inches(11.5), Inches(1.2), size=20, colour=ACCENT)
text(page,
     "Опции «мне не жалко +15 минут, найди дорогу, где я не встану» "
     "сегодня нет ни в одном массовом навигаторе.",
     Inches(0.9), Inches(5.6), Inches(11.5), Inches(1.0), size=18, colour=MUTED)

# ─────────────────────────────────────────────────────────── 3. Idea
page = slide()
heading(page, "Что делает Пробок.Нет")
bullets(page, [
    "Считает долю пути по свободным участкам — а не только ETA.",
    "Сам ведёт поиск: находит пробки на маршруте, запрещает их провайдеру "
    "и просит проложить маршрут заново.",
    "Так по лестнице всё более агрессивных гипотез, пока не кончится бюджет запросов.",
    "Показывает три найденных варианта, отсортированные по доле зелени.",
], Inches(0.9), Inches(2.0), Inches(6.6), Inches(4.0), size=18)
picture(page, "03-route-card-dark.png", Inches(7.8), Inches(2.1), Inches(4.6), Inches(3.2))
text(page, "Карточка варианта: время, расстояние, задержка, доля зелени и полоса "
           "фактических цветов участков.",
     Inches(7.8), Inches(5.5), Inches(4.6), Inches(1.0), size=12, colour=MUTED, italic=True)

# ─────────────────────────────────────────────────────────── 4. Product
page = slide()
heading(page, "Приложение целиком", size=30, underline=False)
picture(page, "01-desktop-dark.png", Inches(0.5), Inches(1.4), W - Inches(1.0), Inches(5.2))
caption(page, "Боевой стек: карта и слой пробок — 2ГИС, маршруты — настоящий ответ провайдера.",
        Inches(6.75))

# ─────────────────────────────────────────────────────────── 5. Top-3
page = slide()
heading(page, "Три варианта, отсортированные по зелени")
picture(page, "04-results-dark.png", Inches(0.9), Inches(1.95), Inches(3.4), Inches(4.9))
bullets(page, [
    "Не «вот маршрут, поверьте» — видно, что поиск действительно был.",
    "Под каждым вариантом полоса фактических цветов участков от провайдера.",
    "Красные и оранжевые минуты показаны отдельной строкой.",
    "Значок ₽ отмечает платную дорогу: порядок в списке определяет только зелень, не цена.",
    "Любой вариант открывается в 2ГИС или Яндексе тем же маршрутом — "
    "через промежуточные точки, а не «примерно туда же».",
], Inches(5.0), Inches(2.1), Inches(7.4), Inches(4.6))

# ─────────────────────────────────────────────────────────── 6. Map
page = slide()
heading(page, "Карта провайдера, маршрут поверх неё", size=30, underline=False)
picture(page, "05-map-traffic-dark.png", Inches(0.5), Inches(1.35), Inches(6.4), Inches(5.3))
bullets(page, [
    "Слой пробок — 2ГИС, вживую.",
    "Цвета участков маршрута берутся у провайдера покусочно.",
    "Нет данных о загруженности — участок UNKNOWN, "
    "зелёным он не считается.",
    "Переключатель провайдера карты: 2ГИС, Яндекс, "
    "OpenStreetMap без ключей.",
], Inches(7.3), Inches(2.2), Inches(5.3), Inches(4.0), size=16)

# ─────────────────────────────────────────────────────────── 7. Search ladder
page = slide()
heading(page, "Как устроен поиск")
steps = [
    ("1", "Кластеризация пробок", "точки RED / ORANGE / YELLOW склеиваются в кластеры радиусом 180–2500 м"),
    ("2", "Защита концов", "кластер ближе 700 м к точке А или Б отбрасывается — иначе провайдер переставит саму точку старта"),
    ("3", "Лестница гипотез", "запрет красных зон → запрет оранжевых → боковой якорь до 25 км → обход двух кластеров → коридор через кольцевую"),
    ("4", "Один запрос на гипотезу", "каждая гипотеза — одно обращение к Routing API"),
    ("5", "Проверка ответа", "концы не уехали дальше 400 м, геометрия не дубль, уложились в лимит км и минут"),
    ("6", "Ранжирование", "сортировка по доле времени по зелёному"),
]
top = Inches(1.95)
for number, title, detail in steps:
    text(page, number, Inches(0.9), top, Inches(0.5), Inches(0.5), size=20, colour=ACCENT, bold=True)
    text(page, title, Inches(1.5), top, Inches(3.2), Inches(0.5), size=17, colour=INK, bold=True)
    text(page, detail, Inches(4.7), top, Inches(7.8), Inches(0.6), size=14, colour=MUTED, spacing=1.15)
    top += Inches(0.78)

# ─────────────────────────────────────────────────────────── 8. Cost
page = slide()
heading(page, "Чего это стоит провайдеру")
stat(page, "6–8", "обращений к Routing API\nна один зелёный поиск", Inches(0.9), Inches(2.2), Inches(3.4))
stat(page, "2 / 4 / 8", "жёсткий бюджет запросов\nпо режимам поиска", Inches(4.7), Inches(2.2), Inches(3.6))
stat(page, "0", "запросов на закладки, избранное\nи всю навигацию по интерфейсу", Inches(8.7), Inches(2.2), Inches(3.8))
text(page,
     ["Два подряд отказа провайдера прекращают поиск досрочно: дневную квоту нельзя "
      "сжигать на заведомо нероутируемые гипотезы.",
      "Сохранённый анализ открывается мгновенно и не тратит ни одного запроса. "
      "Если провайдер недоступен или лимит исчерпан, старые данные не стираются."],
     Inches(0.9), Inches(4.9), Inches(11.5), Inches(1.8), size=17, colour=MUTED)

# ─────────────────────────────────────────────────────────── 9. Modes
page = slide()
heading(page, "Четыре режима — и честный отказ")
rows = [
    ("Режим", "Допуск по расстоянию", "Бюджет запросов"),
    ("Быстрее", "+100 %", "2"),
    ("Баланс", "+100 %", "4"),
    ("Свободнее", "+300 %", "8"),
    ("Только по зелёному", "+300 %", "8"),
]
table = page.shapes.add_table(len(rows), 3, Inches(0.9), Inches(2.0), Inches(7.2), Inches(2.6)).table
table.columns[0].width = Inches(3.0)
table.columns[1].width = Inches(2.4)
table.columns[2].width = Inches(1.8)
for r, row in enumerate(rows):
    for c, value in enumerate(row):
        cell = table.cell(r, c)
        cell.text = value
        cell.fill.solid()
        cell.fill.fore_color.rgb = PAPER_2 if r else RGBColor(0x30, 0x30, 0x36)
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
        run = para.runs[0]
        run.font.size = Pt(14)
        run.font.bold = r == 0
        run.font.name = FONT
        run.font.color.rgb = INK if r else ACCENT
text(page, "Fail-closed", Inches(8.6), Inches(2.05), Inches(4.0), Inches(0.5),
     size=20, colour=ACCENT, bold=True)
text(page,
     "«Только по зелёному»: если ни один найденный вариант не подтверждён зелёным целиком, "
     "приложение не подставляет вместо него обычный быстрый маршрут. Оно показывает три "
     "лучших по зелени и честно помечает, что полную проверку они не прошли.",
     Inches(8.6), Inches(2.7), Inches(4.0), Inches(2.6), size=14, colour=MUTED)
caption(page, "Лимиты объезда задаёт пользователь — до +150 км и +300 мин, — а не зашитая константа.",
        Inches(5.1))

# ─────────────────────────────────────────────────────────── 10. Saved work
page = slide()
heading(page, "Сохранённая работа — без расхода квоты")
picture(page, "07-favourites.png", Inches(0.8), Inches(2.3), Inches(3.7), Inches(1.4))
picture(page, "08-recent.png", Inches(4.8), Inches(2.3), Inches(3.7), Inches(1.4))
picture(page, "09-bookmarks.png", Inches(8.8), Inches(2.3), Inches(3.7), Inches(1.4))
for left, label in ((Inches(0.8), "Избранное"), (Inches(4.8), "Недавние"), (Inches(8.8), "Закладки")):
    text(page, label, left, Inches(3.9), Inches(3.7), Inches(0.4), size=16, colour=ACCENT, bold=True,
         align=PP_ALIGN.CENTER)
bullets(page, [
    "Сохранённый анализ открывается мгновенно и не тратит ни одного запроса к API.",
    "Кнопка обновления рядом: если провайдер недоступен или лимит исчерпан, "
    "старые данные не стираются.",
    "Всё хранится локально в браузере — никакого профиля и никакой отправки маршрутов на сервер.",
], Inches(0.9), Inches(4.6), Inches(11.5), Inches(2.0), size=17)

# ─────────────────────────────────────────────────────────── 11. Devices
page = slide()
heading(page, "Десктоп, планшет, телефон", size=30, underline=False)
picture(page, "12-tablet-light.png", Inches(0.6), Inches(1.5), Inches(7.6), Inches(5.0))
picture(page, "13-mobile-dark.png", Inches(8.6), Inches(1.5), Inches(4.0), Inches(5.0))
caption(page, "Тема запоминается, карта переключает стиль вместе с ней, "
              "на телефоне кнопка поиска помещается на первый экран.", Inches(6.65))

# ─────────────────────────────────────────────────────────── 12. English
page = slide()
heading(page, "Два языка, не только на кнопках", size=30, underline=False)
picture(page, "11-desktop-english.png", Inches(0.5), Inches(1.4), W - Inches(1.0), Inches(5.2))
caption(page, "Английская локаль: интерфейс, режимы, предупреждения и карта.", Inches(6.75))

# ─────────────────────────────────────────────────────────── 13. Architecture
page = slide()
heading(page, "Архитектура")
bullets(page, [
    "edge-api — публичный контракт, OIDC-шлюз, SSE-поток статуса поиска.",
    "routing-orchestrator — сам зелёный поиск: кластеры, гипотезы, бюджет, ранжирование.",
    "provider-yandex — провайдерный слой; второй провайдер подключается "
    "реализацией одного интерфейса.",
    "Веб-приложение на Next.js — планировщик, карта, топ-3, локальные закладки.",
    "Postgres и Redis, OpenTelemetry, метрики Prometheus, дашборды и runbooks.",
], Inches(0.9), Inches(2.0), Inches(11.5), Inches(3.2))
text(page, "Go · Next.js 16 · React · TypeScript · Postgres · Redis · OpenTelemetry · "
           "Docker · Helm · Kubernetes · GitHub Actions",
     Inches(0.9), Inches(5.4), Inches(11.5), Inches(1.0), size=15, colour=MUTED)
text(page, "2GIS Routing API 7.0.0 · 2GIS MapGL · Yandex JS API v3 · OpenStreetMap",
     Inches(0.9), Inches(5.95), Inches(11.5), Inches(1.0), size=15, colour=MUTED)

# ─────────────────────────────────────────────────────────── 14. Readiness
page = slide()
heading(page, "Это не прототип на выходные")
stat(page, "~26 000", "строк кода:\nGo и TypeScript", Inches(0.9), Inches(2.1), Inches(3.0))
stat(page, "85 + e2e", "юнит-тестов и набор\nPlaywright против стека", Inches(4.4), Inches(2.1), Inches(3.4))
stat(page, "14", "ADR — зафиксированные\nархитектурные решения", Inches(8.6), Inches(2.1), Inches(3.6))
bullets(page, [
    "CI: форматирование, линтеры, контрактные тесты, интеграция в docker compose, "
    "e2e в эфемерном окружении, Helm lint и проверка манифестов Kubernetes.",
    "Безопасность: CodeQL, Trivy, gitleaks, аудит зависимостей, SBOM, "
    "keyless-подпись образов, provenance-аттестация.",
    "Поставка: мультиархитектурные образы amd64/arm64 в Docker Hub и GitHub Packages, "
    "Helm-чарт, манифесты Kubernetes, шлюз OIDC.",
], Inches(0.9), Inches(4.4), Inches(11.5), Inches(2.4), size=16)

# ─────────────────────────────────────────────────────────── 15. Demo build
page = slide()
heading(page, "Демо-сборка: без ключей и без запросов", size=30, underline=False)
picture(page, "demo/demo-01-desktop-dark.png", Inches(0.5), Inches(1.4),
        W - Inches(1.0), Inches(5.0))
caption(page, "Статический билд, отвечающий на замороженном настоящем ответе провайдера: "
              "открывается сразу на разобранном примере, всё кликается, ключей не содержит.",
        Inches(6.6))

# ─────────────────────────────────────────────────────────── 16. Fit
page = slide()
heading(page, "Как это встраивается")
options = [
    ("1", "Режим внутри вашего навигатора",
     "«Свободнее» и «Только по зелёному» рядом с «Быстрее» — переключатель, а не отдельное "
     "приложение. Механика готова; нужны покусочная раскраска маршрута и запрос маршрута "
     "с исключёнными зонами."),
    ("2", "Отдельный продукт на вашем API",
     "Приложение работает как есть. Провайдер получает трафик запросов и пользователя, "
     "который платит за качество поездки, а не за минуты."),
    ("3", "Технология целиком",
     "Оркестратор поиска, политика скоринга, контракты и вся обвязка — вместе с историей "
     "решений и тестами."),
]
top = Inches(2.0)
for number, title, detail in options:
    text(page, number, Inches(0.9), top, Inches(0.5), Inches(0.6), size=24, colour=ACCENT, bold=True)
    text(page, title, Inches(1.6), top, Inches(4.4), Inches(0.6), size=18, colour=INK, bold=True)
    text(page, detail, Inches(6.1), top, Inches(6.4), Inches(1.3), size=14, colour=MUTED, spacing=1.2)
    top += Inches(1.5)
text(page, "Первый сценарий интереснее всего провайдеру: покусочные данные о загруженности "
           "уже есть у вас — это то, на чём этот поиск работает и чего нет ни у кого другого.",
     Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.9), size=16, colour=ACCENT)

# ─────────────────────────────────────────────────────────── 17. Limits
page = slide()
heading(page, "Что честно нужно знать", colour=RED)
bullets(page, [
    "Ценность видна в час пик. На свободной дороге зелёного и так 85–90 %, "
    "и разница между вариантами теряется. Продукт нужен в 8:30 и в 18:30.",
    "Поиск стоит запросов: 6–8 вместо одного. При квоте 50 объектов в сутки это шесть "
    "поисков в день — для продукта нужна нормальная квота.",
    "Качество ограничено данными провайдера. Где нет покусочной раскраски — "
    "участок UNKNOWN и в зелень не идёт. Это сделано осознанно и видно в интерфейсе.",
    "Строгая проверка «зелёный целиком» в час пик часто не проходит ни для одного варианта — "
    "и приложение честно об этом говорит.",
], Inches(0.9), Inches(2.1), Inches(11.5), Inches(4.2))

# ─────────────────────────────────────────────────────────── 18. Close
page = slide()
if BRAND.exists():
    page.shapes.add_picture(str(BRAND), Inches(0.9), Inches(1.6), Inches(1.0), Inches(1.0))
text(page, "Работающее приложение, а не идея", Inches(0.9), Inches(2.9), Inches(11), Inches(1.0),
     size=40, colour=INK, bold=True, spacing=1.0)
text(page, "Боевой стек, публичная демо-сборка, полный CI/CD.", Inches(0.9), Inches(4.0),
     Inches(11), Inches(0.7), size=22, colour=ACCENT)
text(page,
     "Проприетарная лицензия: исходный код доступен для чтения, любое использование, "
     "изменение и распространение — только с письменного разрешения автора.",
     Inches(0.9), Inches(4.9), Inches(11), Inches(1.2), size=15, colour=MUTED)

out = PACK / "Пробок.Нет — презентация.pptx"
deck.save(out)
print(f"{out} — {len(deck.slides.__iter__.__self__._sldIdLst)} слайдов, "
      f"{out.stat().st_size // 1024} KB")
